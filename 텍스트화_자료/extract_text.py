#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
자료 모음 폴더의 PDF 및 PPTX 파일을 텍스트로 추출하는 스크립트
작성일: 2025-11-08
"""

import os
import sys
from pathlib import Path
import pdfplumber
from pptx import Presentation
from tqdm import tqdm
import traceback

# 경로 설정
BASE_DIR = Path(__file__).parent.parent
SOURCE_DIR = BASE_DIR / "자료 모음"
OUTPUT_DIR = Path(__file__).parent

# 출력 디렉터리 생성
PDF_OUTPUT_DIR = OUTPUT_DIR / "01_PDF파일"
PPTX_OUTPUT_DIR = OUTPUT_DIR / "02_PPTX파일"
PDF_OUTPUT_DIR.mkdir(exist_ok=True)
PPTX_OUTPUT_DIR.mkdir(exist_ok=True)

def sanitize_filename(filename):
    """파일명에서 특수문자 제거 및 정리"""
    # 확장자 제거
    name = os.path.splitext(filename)[0]
    # 특수문자 제거 및 공백을 언더스코어로 변경
    name = name.replace(" ", "_")
    name = name.replace("(", "").replace(")", "")
    name = name.replace("+", "_")
    name = name.replace(",", "")
    # 연속된 언더스코어 제거
    while "__" in name:
        name = name.replace("__", "_")
    return name[:100]  # 파일명 길이 제한

def extract_pdf_text(pdf_path, output_path):
    """PDF 파일에서 텍스트 추출"""
    try:
        print(f"  처리 중: {pdf_path.name}")
        text_content = []
        text_content.append(f"=" * 80)
        text_content.append(f"파일명: {pdf_path.name}")
        text_content.append(f"파일 크기: {pdf_path.stat().st_size / (1024*1024):.2f} MB")
        text_content.append(f"=" * 80)
        text_content.append("")
        
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
            text_content.append(f"총 페이지 수: {total_pages}\n")
            
            for page_num, page in enumerate(pdf.pages, 1):
                text_content.append(f"\n{'='*80}")
                text_content.append(f"페이지 {page_num}/{total_pages}")
                text_content.append(f"{'='*80}\n")
                
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
                else:
                    text_content.append("[이 페이지는 텍스트를 추출할 수 없습니다. 이미지 기반일 수 있습니다.]")
                
                # 진행상황 표시
                if page_num % 10 == 0:
                    print(f"    진행: {page_num}/{total_pages} 페이지")
        
        # 파일 저장
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))
        
        print(f"  ✅ 완료: {output_path.name} ({total_pages} 페이지)")
        return True, total_pages
        
    except Exception as e:
        error_msg = f"  ❌ 오류 발생: {pdf_path.name}\n     {str(e)}"
        print(error_msg)
        # 오류 로그 저장
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"오류 발생: {pdf_path.name}\n")
            f.write(f"오류 내용: {str(e)}\n\n")
            f.write(traceback.format_exc())
        return False, 0

def extract_pptx_text(pptx_path, output_path):
    """PPTX 파일에서 텍스트 추출"""
    try:
        print(f"  처리 중: {pptx_path.name}")
        text_content = []
        text_content.append(f"=" * 80)
        text_content.append(f"파일명: {pptx_path.name}")
        text_content.append(f"파일 크기: {pptx_path.stat().st_size / (1024*1024):.2f} MB")
        text_content.append(f"=" * 80)
        text_content.append("")
        
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        text_content.append(f"총 슬라이드 수: {total_slides}\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n{'='*80}")
            text_content.append(f"슬라이드 {slide_num}/{total_slides}")
            text_content.append(f"{'='*80}\n")
            
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # 표(table) 처리
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            slide_text.append(row_text)
            
            if slide_text:
                text_content.extend(slide_text)
            else:
                text_content.append("[이 슬라이드는 텍스트가 없습니다. 이미지만 포함되어 있을 수 있습니다.]")
            
            # 진행상황 표시
            if slide_num % 10 == 0:
                print(f"    진행: {slide_num}/{total_slides} 슬라이드")
        
        # 파일 저장
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))
        
        print(f"  ✅ 완료: {output_path.name} ({total_slides} 슬라이드)")
        return True, total_slides
        
    except Exception as e:
        error_msg = f"  ❌ 오류 발생: {pptx_path.name}\n     {str(e)}"
        print(error_msg)
        # 오류 로그 저장
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"오류 발생: {pptx_path.name}\n")
            f.write(f"오류 내용: {str(e)}\n\n")
            f.write(traceback.format_exc())
        return False, 0

def find_files(directory, extension):
    """특정 확장자 파일을 재귀적으로 찾기"""
    files = []
    for root, dirs, filenames in os.walk(directory):
        for filename in filenames:
            if filename.lower().endswith(extension.lower()):
                files.append(Path(root) / filename)
    return files

def main():
    print("\n" + "="*80)
    print("📚 환경과 삶 II - 자료 텍스트화 작업 시작")
    print("="*80 + "\n")
    
    # PDF 파일 처리
    print("\n🔍 PDF 파일 검색 중...")
    pdf_files = find_files(SOURCE_DIR, '.pdf')
    print(f"   발견된 PDF 파일: {len(pdf_files)}개\n")
    
    pdf_success = 0
    pdf_fail = 0
    pdf_total_pages = 0
    
    if pdf_files:
        print("📄 PDF 파일 텍스트 추출 시작\n")
        for pdf_file in pdf_files:
            # 상대 경로 계산 (자료 모음 기준)
            rel_path = pdf_file.relative_to(SOURCE_DIR)
            # 출력 파일명 생성
            safe_name = sanitize_filename(pdf_file.name)
            output_file = PDF_OUTPUT_DIR / f"{safe_name}.txt"
            
            # 이미 처리된 파일은 건너뛰기 (선택사항)
            # if output_file.exists():
            #     print(f"  ⏭️  건너뜀: {pdf_file.name} (이미 처리됨)")
            #     continue
            
            success, pages = extract_pdf_text(pdf_file, output_file)
            if success:
                pdf_success += 1
                pdf_total_pages += pages
            else:
                pdf_fail += 1
        
        print(f"\n✅ PDF 처리 완료: 성공 {pdf_success}개, 실패 {pdf_fail}개")
        print(f"   총 추출 페이지: {pdf_total_pages}페이지\n")
    
    # PPTX 파일 처리
    print("\n🔍 PPTX 파일 검색 중...")
    pptx_files = find_files(SOURCE_DIR, '.pptx')
    print(f"   발견된 PPTX 파일: {len(pptx_files)}개\n")
    
    pptx_success = 0
    pptx_fail = 0
    pptx_total_slides = 0
    
    if pptx_files:
        print("📊 PPTX 파일 텍스트 추출 시작\n")
        for pptx_file in pptx_files:
            # 상대 경로 계산
            rel_path = pptx_file.relative_to(SOURCE_DIR)
            # 출력 파일명 생성
            safe_name = sanitize_filename(pptx_file.name)
            output_file = PPTX_OUTPUT_DIR / f"{safe_name}.txt"
            
            success, slides = extract_pptx_text(pptx_file, output_file)
            if success:
                pptx_success += 1
                pptx_total_slides += slides
            else:
                pptx_fail += 1
        
        print(f"\n✅ PPTX 처리 완료: 성공 {pptx_success}개, 실패 {pptx_fail}개")
        print(f"   총 추출 슬라이드: {pptx_total_slides}개\n")
    
    # HWPX 파일 목록 생성
    print("\n🔍 HWPX 파일 검색 중...")
    hwpx_files = find_files(SOURCE_DIR, '.hwpx')
    print(f"   발견된 HWPX 파일: {len(hwpx_files)}개\n")
    
    if hwpx_files:
        hwpx_list_file = OUTPUT_DIR / "03_HWPX파일_목록.txt"
        with open(hwpx_list_file, 'w', encoding='utf-8') as f:
            f.write("=" * 80 + "\n")
            f.write("HWPX 파일 목록 (수동 변환 필요)\n")
            f.write("=" * 80 + "\n\n")
            f.write(f"총 {len(hwpx_files)}개 파일\n\n")
            
            for i, hwpx_file in enumerate(hwpx_files, 1):
                rel_path = hwpx_file.relative_to(SOURCE_DIR)
                size_mb = hwpx_file.stat().st_size / (1024*1024)
                f.write(f"{i}. {hwpx_file.name}\n")
                f.write(f"   경로: {rel_path}\n")
                f.write(f"   크기: {size_mb:.2f} MB\n\n")
            
            f.write("\n" + "=" * 80 + "\n")
            f.write("처리 방법:\n")
            f.write("1. 한글 프로그램에서 각 파일을 열어 PDF로 저장\n")
            f.write("2. 저장된 PDF 파일을 이 스크립트로 다시 처리\n")
            f.write("3. 또는 온라인 변환 도구 사용 (cloudconvert.com 등)\n")
        
        print(f"⚠️  HWPX 파일 목록 저장: {hwpx_list_file.name}\n")
    
    # 최종 요약
    print("\n" + "="*80)
    print("📊 텍스트화 작업 완료 요약")
    print("="*80)
    print(f"✅ PDF 파일: {pdf_success}/{len(pdf_files)}개 성공 ({pdf_total_pages} 페이지)")
    print(f"✅ PPTX 파일: {pptx_success}/{len(pptx_files)}개 성공 ({pptx_total_slides} 슬라이드)")
    print(f"⚠️  HWPX 파일: {len(hwpx_files)}개 (수동 변환 필요)")
    print(f"\n📁 출력 위치:")
    print(f"   - PDF 텍스트: {PDF_OUTPUT_DIR}")
    print(f"   - PPTX 텍스트: {PPTX_OUTPUT_DIR}")
    print("="*80 + "\n")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n⚠️  사용자에 의해 중단되었습니다.")
        sys.exit(1)
    except Exception as e:
        print(f"\n\n❌ 예상치 못한 오류 발생:")
        print(traceback.format_exc())
        sys.exit(1)
